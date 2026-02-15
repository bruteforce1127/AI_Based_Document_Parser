"""
Document Service - Handles text extraction from multiple file formats
Supports: PDF, Images (OCR), DOCX, PPTX, TXT
"""
import os
import config
import PyPDF2


def extract_text(filepath):
    """
    Extract text from a file based on its extension.
    Returns a list of dicts: [{'page_number': int, 'content': str}, ...]
    """
    ext = os.path.splitext(filepath)[1].lower().lstrip('.')

    extractors = {
        'pdf': extract_text_from_pdf,
        'png': extract_text_from_image,
        'jpg': extract_text_from_image,
        'jpeg': extract_text_from_image,
        'bmp': extract_text_from_image,
        'tiff': extract_text_from_image,
        'webp': extract_text_from_image,
        'docx': extract_text_from_docx,
        'pptx': extract_text_from_pptx,
        'txt': extract_text_from_txt,
    }

    extractor = extractors.get(ext)
    if not extractor:
        print(f"Unsupported file extension: .{ext}")
        return None

    return extractor(filepath)


def extract_text_from_pdf(pdf_path):
    """Extract text content from a PDF file - all pages"""
    pages = []
    try:
        with open(pdf_path, 'rb') as file:
            pdf_reader = PyPDF2.PdfReader(file)
            for page_num in range(len(pdf_reader.pages)):
                page = pdf_reader.pages[page_num]
                text = page.extract_text() or ""
                pages.append({
                    'page_number': page_num + 1,
                    'content': text
                })
    except Exception as e:
        print(f"Error extracting text from PDF: {e}")
        return None
    return pages


def extract_text_from_image(image_path):
    """Extract text from an image using OpenAI Vision API"""
    try:
        import base64
        from PIL import Image
        from services import key_manager
        import io

        # Open and prepare the image
        image = Image.open(image_path)

        # Convert to RGB if necessary (e.g., RGBA PNGs)
        if image.mode in ('RGBA', 'LA', 'P'):
            image = image.convert('RGB')

        # Resize if image is very large (OpenAI has limits)
        max_size = 2048
        if max(image.size) > max_size:
            ratio = max_size / max(image.size)
            new_size = (int(image.size[0] * ratio), int(image.size[1] * ratio))
            image = image.resize(new_size, Image.LANCZOS)

        # Convert image to base64
        buffer = io.BytesIO()
        image.save(buffer, format='JPEG', quality=85)
        base64_image = base64.b64encode(buffer.getvalue()).decode('utf-8')

        # Use OpenAI Vision API to extract text
        client = key_manager.get_openai_client()
        chat_completion = client.chat.completions.create(
            messages=[
                {
                    "role": "user",
                    "content": [
                        {
                            "type": "text",
                            "text": "Extract ALL text visible in this image. Return ONLY the extracted text exactly as it appears, preserving the layout and formatting as much as possible. If there is no text in the image, respond with '[No text detected in image]'."
                        },
                        {
                            "type": "image_url",
                            "image_url": {
                                "url": f"data:image/jpeg;base64,{base64_image}"
                            }
                        }
                    ]
                }
            ],
            model=config.OPENAI_VISION_MODEL,
            temperature=0.1,
            max_tokens=4000
        )

        text = chat_completion.choices[0].message.content.strip()

        if not text or text == '[No text detected in image]':
            return [{'page_number': 1, 'content': '[No text detected in image]'}]

        return [{'page_number': 1, 'content': text}]

    except ImportError:
        print("Pillow not installed. Install with: pip install Pillow")
        return None
    except Exception as e:
        print(f"Error extracting text from image: {e}")
        return None


def extract_text_from_docx(docx_path):
    """Extract text from a DOCX (Word) file"""
    try:
        from docx import Document

        doc = Document(docx_path)
        paragraphs = [p.text for p in doc.paragraphs if p.text.strip()]
        full_text = '\n'.join(paragraphs)

        if not full_text.strip():
            return [{'page_number': 1, 'content': '[No text found in document]'}]

        return [{'page_number': 1, 'content': full_text}]

    except ImportError:
        print("python-docx not installed. Install with: pip install python-docx")
        return None
    except Exception as e:
        print(f"Error extracting text from DOCX: {e}")
        return None


def extract_text_from_pptx(pptx_path):
    """Extract text from a PPTX (PowerPoint) file - one page per slide"""
    try:
        from pptx import Presentation

        prs = Presentation(pptx_path)
        pages = []

        for slide_num, slide in enumerate(prs.slides, 1):
            texts = []
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for paragraph in shape.text_frame.paragraphs:
                        text = paragraph.text.strip()
                        if text:
                            texts.append(text)

            slide_content = '\n'.join(texts) if texts else '[No text on this slide]'
            pages.append({
                'page_number': slide_num,
                'content': slide_content
            })

        return pages if pages else [{'page_number': 1, 'content': '[Empty presentation]'}]

    except ImportError:
        print("python-pptx not installed. Install with: pip install python-pptx")
        return None
    except Exception as e:
        print(f"Error extracting text from PPTX: {e}")
        return None


def extract_text_from_txt(txt_path):
    """Extract text from a plain text file"""
    try:
        # Try UTF-8 first, then fall back to latin-1
        try:
            with open(txt_path, 'r', encoding='utf-8') as f:
                text = f.read()
        except UnicodeDecodeError:
            with open(txt_path, 'r', encoding='latin-1') as f:
                text = f.read()

        if not text.strip():
            return [{'page_number': 1, 'content': '[Empty file]'}]

        return [{'page_number': 1, 'content': text}]

    except Exception as e:
        print(f"Error reading text file: {e}")
        return None


def get_full_text(pages):
    """Combine all pages into full text"""
    if not pages:
        return ""
    return '\n\n'.join([f"--- Page {p['page_number']} ---\n{p['content']}" for p in pages])
